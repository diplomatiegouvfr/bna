from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

import powerpoint 
from calcul import frq_mot
from utils import format_nbr
        
from datetime import date
import locale

import os
import os.path


def generate_ppt(short_name_file, file_without_extension, categories_dict, short_name_file_dico_freq, most_frq_word):
    # Creating Object 
    path = 'Lib/BNA.pptx'
    prs = Presentation(path)

    # Pour obtenir la date en français
    locale.setlocale(locale.LC_ALL, 'fr_FR.utf8')

    today = date.today()
    
    # dd/mm/YY
    d1 = today.strftime("%d/%m/%Y")
    annee = today.strftime("%Y")
    
    # slides layouts
    slide_pres = prs.slide_layouts[1]
    slide_info = prs.slide_layouts[2]
    slide_table = prs.slide_layouts[6]
    slide_big_img = prs.slide_layouts[7]
    slide_medium_img = prs.slide_layouts[5]
    slide_categ = prs.slide_layouts[4]
    slide_wc = prs.slide_layouts[3]
    
    
    # Slide 1
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide1 = prs.slides.add_slide(slide_pres)
    
    for shape in slide1.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
    
    text_frame.text = f"Baromètre BNA {annee} \nanalyses et détails […]".upper()
    
    p = text_frame.add_paragraph()
    p.text = f"\nVersion du {today.strftime('%d %B %Y')}".upper()
    p.level = 1
    p.font.size = Pt(18) 
    
    # Slide 2
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide2 = prs.slides.add_slide(slide_info)
    
    title2 = slide2.shapes.title
    title2.text = f"Sur le baromètre BNA {annee}"
    title2.size = Pt(28)
    
    for shape in slide2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
    
    text_frame.text = f"À propos du baromètre BNA {annee}"
    p = text_frame.add_paragraph()
    p.text = "  Organisé par la DiNum et traité par l’institut OpinionWay"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = f"  Durant la période du 1 au 23 juin {annee}"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Avec la participation de 614 137 agents issus de 13 ministères"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Taux moyen de participation de 18% pour les ministères"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "\nSur la participation [...]"
    
    p = text_frame.add_paragraph()
    p.text = "  [...] agents contactés sur la période (soit [...]% du global)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Taux de participation moyen de [...]% (soit [...] agents)"
    p.level = 1
    
    # Nombre de réponses fournies
    file = open(f"entrees/{short_name_file}.txt", "rt", encoding="utf-8") # fichier d'extraction de tous les mots        si erreur : , encoding="utf-8", errors="ignore"
    data = file.read()
    row_count = data.split("\n")
    reponses_fournies = len(row_count) - 1
    
    # Nombre de mots pris en compte
    file = open(f"entrees/{short_name_file}.txt", "rt", encoding="utf-8") # fichier d'extraction de tous les mots        si erreur : , encoding="utf-8", errors="ignore"
    data = file.read()
    words = data.split()
    mots = len(words)
    
    # Nombre d'expressions retenues
    file = open(f'sorties/{short_name_file}/{file_without_extension}.txt', 'rt', encoding='utf-8') # fichier en sortie du premier traitement
    data = file.read()
    lines = data.split("\n")
    expressions = len(lines) - 1
    
    # Nombre de mots pris en compte
    file = open(f'sorties/{short_name_file}/{file_without_extension}.txt', 'rt', encoding='utf-8') # fichier en sortie du premier traitement
    data = file.read()
    words = data.split()
    mots1 = len(words)
    
    # Nombre de couples retenus
    file = open(f'sorties_v4/{file_without_extension}/{file_without_extension}_couples.txt', 'rt', encoding='utf-8') # fichier en sortie du deuxieme traitement
    data = file.read()
    lines = data.split("\n")
    couples = len(lines) - 1
    
    # Nombre de mots pris en compte
    file = open(f'sorties_v4/{file_without_extension}/{file_without_extension}_couples.txt', 'rt', encoding='utf-8') # fichier en sortie du deuxieme traitement
    data = file.read()
    data = [i.replace('\n', ' ') for i in data]
    data1 = ''.join(data)
    words = data1.split()
    mots2 = len(words)
    
    # Ratio 1 : des expressions retenues
    ratio1 = round((expressions/reponses_fournies)*100)
    
    # Ratio 2 : des mots pris en compte
    ratio2 = round((mots1/mots)*100)
    
    # Ratio 3 : des couples retenus
    ratio3 = round((couples/expressions)*100)
    
    # Ratio 4 : des mots pris en compte
    ratio4 = round((mots2/mots1)*100)
    
    def set_superscript(font):
        font._element.set('baseline', '30000')
    
    
    p = text_frame.add_paragraph()
    p.text = "\nSur les données traitées"
    
    p = text_frame.add_paragraph()
    p.text = f"  Analyse des {format_nbr(reponses_fournies)} réponses fournies, dont {format_nbr(mots)} mots pris en compte" #format_nbr()
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = f"  1er traitement, {format_nbr(expressions)} expressions retenues (~{format_nbr(ratio1)}%), dont {format_nbr(mots1)} mots pris en compte (~{format_nbr(ratio2)}%)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = f"  2ème traitement, {format_nbr(couples)} couples retenus (~{format_nbr(ratio3)}%), dont {format_nbr(mots2)} mots pris en compte (~{format_nbr(ratio4)}%)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  La faible quantité de réponses a limité le traitement à certaines rubriques uniquement"
    p.level = 1
    
    
    # Slide 3
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide3 = prs.slides.add_slide(slide_info)
    
    title3 = slide3.shapes.title
    title3.text = "Sur l’analyse des données"
    
    for shape in slide3.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame
    
    text_frame.text = "Objectifs de l’étude"
    
    p = text_frame.add_paragraph()
    p.text = "  Étudier les réponses aux questions ouvertes"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Fournir des analyses sur les expressions utilisées"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Établir les tendances et graphiques associés"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "\nSur la méthodologie utilisée au MEAE"
    
    p = text_frame.add_paragraph()
    p.text = "  Indexation des questions ouvertes (Q12 et Q30)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Analyse des phrases pour reconnaître les négations, les formes actives ou passives etc."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Hiérarchisation des réponses par grandes familles pour définir les tendances (grammaire de Chomsky)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Construction des liens entre les expressions retenues (vectorisation de Mikolov)"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Corrélation des 3 résultats (fichier pivot, 9 fichiers catégories/mots-clés et analyse fréquentielle) sur 5 points principaux possédant potentiellement le meilleur ratio amélioration/ressenti"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Recoupement avec analyse fréquentielle des données brutes initiales"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Analyse des fréquences des expressions et graphiques"
    p.level = 1
    
    # Slide 4
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide4 = prs.slides.add_slide(slide_table)
    
    title4 = slide4.shapes.title
    title4.text = "Sur les expressions recoupées"
    title4.size = Pt(28)
    
    placeholder = slide4.placeholders[11]
    table1_frame = placeholder.insert_table(22,15) #, left_table, top_table, width_table, height_table
    
    table1 = table1_frame.table
    
    tbl = table1._graphic_frame.element.graphic.graphicData.tbl
    style_id ='{2D5ABB26-0587-4C30-8999-92F81FD0307C}' # See github for UUID
    tbl[0][-1].text = style_id
    
    # set title and width
    for p in range(0,15,3):
        cell = table1.cell(0,p)
        cell.text = 'Rang'
        table1.columns[p].width = Inches(0.5)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
    for p in range(1,15,3):
        cell = table1.cell(0,p)
        cell.text = 'Freq'
        table1.columns[p].width = Inches(0.5)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
    for p in range(2,15,3):
        cell = table1.cell(0,p)
        cell.text = 'Mot'
        table1.columns[p].width = Inches(1.2)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
    
    # including ranking
    b = 0
    for p in range(0,15,3):
        for i in range(1,20):
            b = b + 1
            cell = table1.cell(i,p)
            cell.text = str(b)
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            
    # including freq
    a = short_name_file_dico_freq # a = frq_mot_dico("entrees/{}".format(short_name_file))
    b = 0
    for p in range(1,15,3):
        for i in range(1,20):
            cell = table1.cell(i,p)
            cell.text = [str(value) for value in a.values()][b]
            b = b + 1
    
    # including words
    b = 0
    for p in range(2,15,3):
        for i in range(1,20):
            cell = table1.cell(i,p)
            cell.text = [key for key in a.keys()][b]
            b = b + 1
    
    # formatting all cells
    for i in range(20): # loop all lines
        for p in range(15): # loop all colums
            cell = table1.cell(i,p)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(8)
            para.font.name = 'Arial'
    
    for i in range(20): # loop all lines
        table1.rows[i].height = Inches(0.01)
    
    powerpoint.slide_text(slide4, '\nLes différentes expressions retenues sont classées par ordre d’occurrence dans l’ensemble des réponses ouvertes fournies par le baromètre. ')
    
    
    # Slide 5
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide5 = prs.slides.add_slide(slide_big_img)
    
    title5 = slide5.shapes.title
    title5.text = "Sur la perception de l’équipement principal"
    
    placeholder_picture = slide5.shapes[1]
    placeholder_picture.insert_picture(rf"sorties_v4/{file_without_extension}/images/wordcloud_output_principal.png")
    
    
    expression = short_name_file_dico_freq # frq_mot_dico("entrees/{}".format(short_name_file))
    expression = list(expression.keys())[0]
    
    placeholder_text = slide5.placeholders[11]
    placeholder_text.text = f"\nLes différentes expressions sont classées par ordre d’occurrence dans l’ensemble des réponses ouvertes du baromètre. La taille des caractères donne le niveau de citation ; ainsi « {expression} » est l’expression la plus citée pour cette question."
    
    
    # Slide 6
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide6 = prs.slides.add_slide(slide_medium_img)
    
    title6 = slide6.shapes.title
    title6.text = "Sur l'analyse du verbatim en 9 grandes familles"
    
    # adding a picture
    placeholder = slide6.shapes[1]  # idx key, not position
    placeholder.insert_picture(rf'sorties_v4/{file_without_extension}/images/ExportedPieChart.png')
    '''
    for shape in slide6.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame'
    '''
    
    placeholder_text = slide6.placeholders[11]
    placeholder_text.text = '\nLes expressions retenues des questions ouvertes sont hiérarchisées en 9 catégories en fonction du lien entre elles pour la question posée et sont analysées par pertinence. Ainsi, [...]% des réponses sont comprises entre "très satisfaisant" et "assez passable".'
    
    
    # Slide 7
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide7 = prs.slides.add_slide(slide_medium_img)
    
    title7 = slide7.shapes.title
    title7.text = "Sur l'analyse du verbatim en 9 grandes familles"
    
    # adding a picture
    placeholder = slide7.shapes[1]  # idx key, not position
    placeholder.insert_picture(rf'sorties_v4/{file_without_extension}/images/gauss_function_stat.png')
    
    placeholder_text = slide7.placeholders[11]
    placeholder_text.text = '\nLes expressions retenues des questions ouvertes sont hiérarchisées en fonction du lien entre elles pour la question posée et sont analysées par rapport à la distribution moyenne attendue. Idéalement la dispersion des valeurs observée doit être proche de zéro pourcent, donc la rubrique « passable » est trop représentée.'
    
    
    # Slide 8
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide8 = prs.slides.add_slide(slide_categ)
    
    
    title8 = slide8.shapes.title
    title8.text = "Sur l’analyse des données pour satisfaisant"
    title8.size = Pt(28)
    
    placeholder_picture = slide8.placeholders[10]
    placeholder_text = slide8.placeholders[17]
    powerpoint.check_picture('assez_satisfaisant', placeholder_text, placeholder_picture, file_without_extension)
    
    placeholder_text = slide8.shapes[2]
    placeholder_text.text = '\nLes expressions sont classées en fonction de l’appréciation faite. Dans cette partie de l’analyse, les expressions sont traitées par rapport au critère « satisfaisant » et ses dérivées. Les planches suivantes détaillent ces données. '
    
    placeholder_text2 = slide8.shapes[3]  # idx key, not position
    placeholder_text2.text = '    Assez satisfaisant         	      Satisfaisant                  Très satisfaisant'
    
    
    placeholder_picture = slide8.placeholders[15]
    placeholder_text = slide8.placeholders[18]
    powerpoint.check_picture('satisfaisant', placeholder_text, placeholder_picture, file_without_extension)
    
    placeholder_picture = slide8.placeholders[16]
    placeholder_text = slide8.placeholders[19]
    powerpoint.check_picture('tres_satisfaisant', placeholder_text, placeholder_picture, file_without_extension)
    
    
    # Slide 9
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide9 = prs.slides.add_slide(slide_wc)
    
    title9 = slide9.shapes.title
    title9.text = "Sur l’analyse des données pour assez satisfaisant"
    title9.size = Pt(28)
    
    placeholder_picture = slide9.placeholders[10]
    placeholder_text = slide9.placeholders[12]
    pic = powerpoint.check_picture('assez_satisfaisant', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide9.placeholders[11]
    powerpoint.commentaires('assez_satisfaisant', placeholder_text2, pic, frq_mot(categories_dict.get('assez_satisfaisant')), most_frq_word.get('assez_satisfaisant'))
    
    
    # Slide 10
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide10 = prs.slides.add_slide(slide_wc)
    
    title10 = slide10.shapes.title
    title10.text = "Sur l’analyse des données pour satisfaisant"
    title10.size = Pt(28)
    
    placeholder_picture = slide10.shapes[1]
    placeholder_text = slide10.placeholders[12]
    pic = powerpoint.check_picture('satisfaisant', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide10.shapes[2]
    powerpoint.commentaires('satisfaisant', placeholder_text2, pic, frq_mot(categories_dict.get('satisfaisant')), most_frq_word.get('satisfaisant'))
    
    
    # Slide 11
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide11 = prs.slides.add_slide(slide_wc)
    
    title11 = slide11.shapes.title
    title11.text = "Sur l’analyse des données pour très satisfaisant"
    title11.size = Pt(28)
    
    placeholder_picture = slide11.shapes[1]
    placeholder_text = slide11.placeholders[12]
    pic = powerpoint.check_picture('tres_satisfaisant', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide11.shapes[2]
    powerpoint.commentaires('tres_satisfaisant', placeholder_text2, pic, frq_mot(categories_dict.get('tres_satisfaisant')), most_frq_word.get('tres_satisfaisant'))
    
    
    # Slide 12
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide12 = prs.slides.add_slide(slide_categ)
    
    title12 = slide12.shapes.title
    title12.text = "Sur l’analyse des données pour passable"
    title12.size = Pt(28)
    
    placeholder_picture = slide12.placeholders[10]
    placeholder_text = slide12.placeholders[17]
    powerpoint.check_picture('assez_passable', placeholder_text, placeholder_picture, file_without_extension)
    
    placeholder_text1 = slide12.shapes[2]
    placeholder_text1.text = '\nLes expressions sont classées en fonction de l’appréciation faite. Dans cette partie de l’analyse, les expressions sont traitées par rapport au critère « passable » et ses dérivées. Les planches suivantes détaillent ces données.'
    
    placeholder_text2 = slide12.shapes[3]
    placeholder_text2.text = '      Assez passable                       Passable                      Très passable'
    
    
    placeholder_picture = slide12.placeholders[15]
    placeholder_text = slide12.placeholders[18]
    powerpoint.check_picture('passable', placeholder_text, placeholder_picture, file_without_extension)
    
    placeholder_picture = slide12.placeholders[16]
    placeholder_text = slide12.placeholders[19]
    powerpoint.check_picture('tres_passable', placeholder_text, placeholder_picture, file_without_extension)
    
    
    # Slide 13
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide13 = prs.slides.add_slide(slide_wc)
    
    title13 = slide13.shapes.title
    title13.text = "Sur l’analyse des données pour assez passable"
    title13.size = Pt(28)
    
    placeholder_picture = slide13.shapes[1]
    placeholder_text = slide13.placeholders[12]
    pic = powerpoint.check_picture('assez_passable', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide13.shapes[2]
    powerpoint.commentaires('assez_passable', placeholder_text2, pic, frq_mot(categories_dict.get('assez_passable')), most_frq_word.get('assez_passable'))
    
    
    # Slide 14
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide14 = prs.slides.add_slide(slide_wc)
    
    title14 = slide14.shapes.title
    title14.text = "Sur l’analyse des données pour passable"
    title14.size = Pt(28)
    
    placeholder_picture = slide14.shapes[1]
    placeholder_text = slide14.placeholders[12]
    pic = powerpoint.check_picture('passable', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide14.shapes[2]
    powerpoint.commentaires('passable', placeholder_text2, pic, frq_mot(categories_dict.get('passable')), most_frq_word.get('passable'))
    
    
    # Slide 15
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide15 = prs.slides.add_slide(slide_wc)
    
    title15 = slide15.shapes.title
    title15.text = "Sur l’analyse des données pour très passable"
    title15.size = Pt(28)
    
    placeholder_picture = slide15.shapes[1]
    placeholder_text = slide15.placeholders[12]
    pic = powerpoint.check_picture('tres_passable', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide15.shapes[2]
    powerpoint.commentaires('tres_passable', placeholder_text2, pic, frq_mot(categories_dict.get('tres_passable')), most_frq_word.get('tres_passable'))
    
    
    # Slide 16
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide16 = prs.slides.add_slide(slide_categ)
    
    title16 = slide16.shapes.title
    title16.text = "Sur l’analyse des données pour insatisfaisant"
    title16.size = Pt(28)
    
    placeholder_picture = slide16.placeholders[10]
    placeholder_text = slide16.placeholders[17]
    powerpoint.check_picture('assez_insatisfaisant', placeholder_text, placeholder_picture, file_without_extension)
    
    placeholder_text1 = slide16.shapes[2]
    placeholder_text1.text = '\nLes expressions sont classées en fonction de l’appréciation faite. Dans cette partie de l’analyse, les expressions sont traitées par rapport au critère « insatisfaisant » et ses dérivées. Les planches suivantes détaillent ces données.'
    
    placeholder_text2 = slide16.shapes[3]
    placeholder_text2.text = '  Assez insatisfaisant        	     Insatisfaisant                Très insatisfaisant'
    
    placeholder_picture = slide16.placeholders[15]
    placeholder_text = slide16.placeholders[18]
    powerpoint.check_picture('insatisfaisant', placeholder_text, placeholder_picture, file_without_extension)
    
    placeholder_picture = slide16.placeholders[16]
    placeholder_text = slide16.placeholders[19]
    powerpoint.check_picture('tres_insatisfaisant', placeholder_text, placeholder_picture, file_without_extension)
    
    
    # Slide 17
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide17 = prs.slides.add_slide(slide_wc)
    
    title17 = slide17.shapes.title
    title17.text = "Sur l’analyse des données pour assez insatisfaisant"
    title17.size = Pt(28)
    
    placeholder_picture = slide17.shapes[1]
    placeholder_text = slide17.placeholders[12]
    pic = powerpoint.check_picture('assez_insatisfaisant', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide17.shapes[2]
    powerpoint.commentaires('assez_insatisfaisant', placeholder_text2, pic, frq_mot(categories_dict.get('assez_insatisfaisant')), most_frq_word.get('assez_insatisfaisant'))
    
    
    # Slide 18
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide18 = prs.slides.add_slide(slide_wc)
    
    title18 = slide18.shapes.title
    title18.text = "Sur l’analyse des données pour insatisfaisant"
    title18.size = Pt(28)
    
    placeholder_picture = slide18.shapes[1]
    placeholder_text = slide18.placeholders[12]
    pic = powerpoint.check_picture('insatisfaisant', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide18.shapes[2]
    powerpoint.commentaires('insatisfaisant', placeholder_text2, pic, frq_mot(categories_dict.get('insatisfaisant')), most_frq_word.get('insatisfaisant'))
    
    
    # Slide 19
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide19 = prs.slides.add_slide(slide_wc)
    
    title19 = slide19.shapes.title
    title19.text = "Sur l’analyse des données pour très insatisfaisant"
    title19.size = Pt(28)
    
    placeholder_picture = slide19.shapes[1]
    placeholder_text = slide19.placeholders[12]
    pic = powerpoint.check_picture('tres_insatisfaisant', placeholder_text, placeholder_picture, file_without_extension)[1]
    
    placeholder_text2 = slide19.shapes[2]
    powerpoint.commentaires('tres_insatisfaisant', placeholder_text2, pic, frq_mot(categories_dict.get('tres_insatisfaisant')), most_frq_word.get('tres_insatisfaisant'))
    
    
    # Slide 20
    #-----------------------------------------------------------------------------------------------------------------------
    slide20 = prs.slides.add_slide(slide_table)
    
    title20 = slide20.shapes.title
    title20.text = "Les suggestions d’améliorations"
    
    #Add table to Slide
    placeholder = slide20.placeholders[11]
    table1_frame = placeholder.insert_table(6,3)
    
    #Populating a table
    table1 = table1_frame.table
    
    cell = table1.cell(0,0)
    cell.text = "Mot-clef"
    
    cell = table1.cell(0,1)
    cell.text = "Idée principale"
    
    cell = table1.cell(0,2)
    cell.text = "Réponse-type au baromètre"
    
    # including words
    b = 0
    for i in range(1,6):
        cell = table1.cell(i,0)
        cell.text = [key for key in a.keys()][b]
        b = b + 1
        
    # including main ideas
    for i in range(1,6):
        cell = table1.cell(i,1)
        cell.text = " Texte à modifier"
    
    # including sentences
    b = 0
    for i in range(1,6):
        cell = table1.cell(i,2)
        mot = [key for key in a.keys()][b]
        cell.text = str(powerpoint.reponse(short_name_file, mot, placeholder_text))
        b = b + 1
    
    #Applaying table styles
    tbl = table1._graphic_frame.element.graphic.graphicData.tbl
    style_id ='{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}' # See github for UUID
    tbl[0][-1].text = style_id

    # formatting all cells
    for i in range(6): # loop all lines
            for p in range(3): # loop all colums
                cell = table1.cell(i,p)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(16)
                para.font.name = 'Arial'
    
    # formatting third column cells of second row
    for i in range(5):
        cell = table1.cell(i+1,2) # From second line of third column to fifth line of third column
        para = cell.text_frame.paragraphs[0]
        para.font.italic = True
        para.font.size = Pt(14)
        para.font.name = 'Arial'
    
    # formatting third column cells of second row
    for i in range(5):
        for p in range(3):
            cell = table1.cell(i+1,p) # From second line of all columns to fifth line
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT # left align
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE # center align on cell
    
    # set column widths
    table1.columns[0].width = Inches(1.5)
    table1.columns[1].width = Inches(2.5)
    table1.columns[2].width = Inches(6.7) # sum must be equal to 10.63
    
    
    # Slide 21
    #-----------------------------------------------------------------------------------------------------------------------
    
    slide21 = prs.slides.add_slide(slide_info)
    
    title21 = slide21.shapes.title
    title21.text = "Les biais du traitement effectué"
    title21.size = Pt(28)
    
    for shape in slide21.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame
    
    text_frame.text = "Traitement des données [...]"
    
    p = text_frame.add_paragraph()
    p.text = "Comme tout exercice de ce type, les résultats sont souvent imparfaits. Il est nécessaire de ne pas conclure trop hâtivement sans avoir mis au préalable en corrélation les conclusions des représentations graphiques faites avec les données traitées et les fichiers intermédiaires produits :\n"
    p.font.size = Pt(16) 
    
    p = text_frame.add_paragraph()
    p.text = "  Biais 1 - certains logiciels spécifiques à chaque ministère apparaissent sans que le traitement en connaisse le contexte spécifique (Diplomatie, Sphère, Scénarios, etc.) ou puisse avoir les ressources nécessaires afin de l’interpréter avec justesse\n"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Biais 2 - certaines thématiques parcourent pratiquement l’intégralité des 9 indices de satisfaction, par exemple le « stockage » (souvent en rapport avec la taille des messageries). Ceci est normal et exprime la diversité des perceptions des utilisateurs sur les thématiques retenues en fonction des repères propres à chacun"
    p.level = 1
    
    
    # Slide 22
    #-----------------------------------------------------------------------------------------------------------------------

    slide22 = prs.slides.add_slide(slide_info)
    
    title22 = slide22.shapes.title
    title22.text = "Conclusion et limites de l’analyse"
    
    for shape in slide22.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
    
    text_frame.text = "Traitement des données du ministère"
    
    p = text_frame.add_paragraph()
    p.text = "  Analyse faite sur les données brutes et hors-contexte"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Certaines expressions ou conclusions élaborées sont à interpréter"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Certaines expressions sont des appellations internes à traiter par le ministère"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "\nSur la participation du ministère"
    
    p = text_frame.add_paragraph()
    p.text = "  Travail effectué uniquement sur les données transmises pour le baromètre"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Méthodologie simple et fonctionnelle déjà testée en interne"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "  Des annexes et compléments sont disponibles"
    p.level = 1    
  
    
    # Dates 
    placeholder_list = [
        slide2.placeholders[12],
        slide3.placeholders[12],
        slide4.placeholders[13],
        slide5.placeholders[12],
        slide6.placeholders[12],
        slide7.placeholders[12],
        slide8.placeholders[12],
        slide9.placeholders[13],
        slide10.placeholders[13],
        slide11.placeholders[13],
        slide12.placeholders[12],
        slide13.placeholders[13],
        slide14.placeholders[13],
        slide15.placeholders[13],
        slide16.placeholders[12],
        slide17.placeholders[13],
        slide18.placeholders[13],
        slide19.placeholders[13],
        slide20.placeholders[13],
        slide21.placeholders[12],
        slide22.placeholders[12]
        ]
    
    for placeholder in placeholder_list:
        placeholder.text = d1
    
    # set the proofing language
    new_language = MSO_LANGUAGE_ID.FRENCH
    
    # apply the proofing language
    # iterate through all slides
    for slide_no, slide in enumerate(prs.slides):
        # iterate through all shapes/objects on one slide
        for shape in slide.shapes:
            # check if the shape/object has text (pictures e.g. don't have text)
            if shape.has_text_frame:
                # check for each paragraph of text for the actual shape/object
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # set the 'new_language'
                        run.font.language_id = new_language
            # check if the shape/object has table
            if shape.has_table:
                tbl = shape.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0, row_count):
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        paragraphs = cell.text_frame.paragraphs 
                        # check for each paragraph of text for the actual cell
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                # set the 'new_language
                                run.font.language_id = new_language
            else:
                pass
    
    
    # saving file
    barometre = short_name_file
    prs.save(f'{barometre}.pptx')
    os.startfile(f'{barometre}.pptx')
    
    return print(f'{short_name_file}.pptx créé')







